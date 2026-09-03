"""Unit tests for supported knowledge-document parsers."""

from __future__ import annotations

from pathlib import Path

import pytest

from octop.infra.knowledge.parse import parse_document


def test_parse_plain_text_and_markdown(tmp_path: Path) -> None:
    text = tmp_path / "notes.txt"
    markdown = tmp_path / "readme.md"
    text.write_text("plain notes", encoding="utf-8")
    markdown.write_text("# Heading\n\nbody", encoding="utf-8")

    assert parse_document(text) == "plain notes"
    assert parse_document(markdown) == "# Heading\n\nbody"


def test_parse_pdf_docx_and_pptx(tmp_path: Path) -> None:
    from docx import Document
    from pptx import Presentation
    from pypdf import PdfWriter

    pdf = tmp_path / "empty.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=72, height=72)
    with pdf.open("wb") as output:
        writer.write(output)

    docx = tmp_path / "notes.docx"
    document = Document()
    document.add_paragraph("Word notes")
    document.save(docx)

    pptx = tmp_path / "slides.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[0]).shapes.title.text = "Slide title"
    presentation.save(pptx)

    assert parse_document(pdf) == ""
    assert parse_document(docx) == "Word notes"
    assert parse_document(pptx) == "Slide title"


def test_parse_csv_xlsx_and_xls(tmp_path: Path) -> None:
    import xlwt
    from openpyxl import Workbook

    csv_path = tmp_path / "sales.csv"
    csv_path.write_text("item,qty\napple,2\n", encoding="utf-8")

    tsv_path = tmp_path / "sales.tsv"
    tsv_path.write_text("item\tqty\napple\t2\n", encoding="utf-8")

    xlsx = tmp_path / "sales.xlsx"
    workbook = Workbook()
    sheet = workbook.active
    assert sheet is not None
    sheet.title = "Q1"
    sheet.append(["item", "qty"])
    sheet.append(["apple", 2])
    extra = workbook.create_sheet("Q2")
    extra.append(["item", "qty"])
    extra.append(["pear", 3])
    workbook.save(xlsx)

    xlsm = tmp_path / "sales.xlsm"
    workbook.save(xlsm)

    xls = tmp_path / "sales.xls"
    book = xlwt.Workbook()
    q1 = book.add_sheet("Q1")
    q1.write(0, 0, "item")
    q1.write(0, 1, "qty")
    q1.write(1, 0, "apple")
    q1.write(1, 1, 2)
    book.save(str(xls))

    expected_xlsx = "# Q1\nitem\tqty\napple\t2\n\n# Q2\nitem\tqty\npear\t3"
    assert parse_document(csv_path) == "# sales\nitem\tqty\napple\t2"
    assert parse_document(tsv_path) == "# sales\nitem\tqty\napple\t2"
    assert parse_document(xlsx) == expected_xlsx
    assert parse_document(xlsm) == expected_xlsx
    assert parse_document(xls) == "# Q1\nitem\tqty\napple\t2"


def test_parse_html_json_and_plain_variants(tmp_path: Path) -> None:
    html = tmp_path / "page.html"
    html.write_text(
        "<html><head><style>p{color:red}</style></head>"
        "<body><h1>Title</h1><p>Hello <b>world</b></p>"
        "<script>alert(1)</script></body></html>",
        encoding="utf-8",
    )
    markdown = tmp_path / "notes.markdown"
    markdown.write_text("# Heading\n\nbody", encoding="utf-8")
    rst = tmp_path / "notes.rst"
    rst.write_text("Heading\n=======\n\nbody", encoding="utf-8")
    yaml_path = tmp_path / "config.yaml"
    yaml_path.write_text("name: octop\n", encoding="utf-8")
    json_path = tmp_path / "data.json"
    json_path.write_text('{"name":"octop","ok":true}', encoding="utf-8")
    jsonl = tmp_path / "rows.jsonl"
    jsonl.write_text('{"a":1}\n{"b":2}\n', encoding="utf-8")

    assert parse_document(html) == "Title\nHello world"
    assert parse_document(markdown) == "# Heading\n\nbody"
    assert parse_document(rst) == "Heading\n=======\n\nbody"
    assert parse_document(yaml_path) == "name: octop\n"
    assert parse_document(json_path) == '{\n  "name": "octop",\n  "ok": true\n}'
    assert parse_document(jsonl) == '{"a":1}\n{"b":2}\n'


def test_parse_rejects_unsupported_extension(tmp_path: Path) -> None:
    path = tmp_path / "unsupported.zip"
    path.write_bytes(b"not a document")

    with pytest.raises(ValueError, match="unsupported"):
        parse_document(path)
